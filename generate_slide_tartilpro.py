#!/usr/bin/env python3
"""
Generate TartilPro x SD Khadijah 3 slide deck.
Theme: Modern Islamic Tech (Emerald Green, Soft Gold, White).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Theme colors
EMERALD = RGBColor(4, 120, 87)          # #047857
EMERALD_DARK = RGBColor(6, 78, 59)      # #064E3B
GOLD = RGBColor(212, 175, 55)           # #D4AF37
GOLD_LIGHT = RGBColor(250, 240, 210)    # #FAF0D2
WHITE = RGBColor(255, 255, 255)         # #FFFFFF
SOFT_BG = RGBColor(236, 253, 245)       # #ECFDF5
DARK_TEXT = RGBColor(31, 41, 55)        # #1F2937
MUTED_TEXT = RGBColor(75, 85, 99)       # #4B5563

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"


def add_full_shape(slide, shape_type, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0)  # no border
    else:
        shape.line.fill.background()
    return shape


def set_text_frame(text_frame, text, font_name, size, color, bold=False, align=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return p


def add_text_box(slide, left, top, width, height, text, font_name, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    set_text_frame(tf, text, font_name, size, color, bold=bold, align=align)
    return box


def add_bullet_list(slide, left, top, width, height, items, font_name, size, color, bullet_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = font_name
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
    return box


def add_footer(slide):
    # Gold accent line at bottom
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_HEIGHT - Inches(0.35),
        SLIDE_WIDTH, Inches(0.08),
        fill=GOLD
    )
    # Footer text
    add_text_box(
        slide, Inches(0.5), SLIDE_HEIGHT - Inches(0.35),
        Inches(6), Inches(0.3),
        "TartilPro x SD Khadijah 3",
        FONT_BODY, Pt(12), WHITE, bold=True
    )
    add_text_box(
        slide, SLIDE_WIDTH - Inches(3.5), SLIDE_HEIGHT - Inches(0.35),
        Inches(3), Inches(0.3),
        "Presentasi Produk | 2026",
        FONT_BODY, Pt(12), WHITE, align=PP_ALIGN.RIGHT
    )


def add_header_bar(slide, title, subtitle=None):
    # Emerald header bar
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.15),
        fill=EMERALD
    )
    # Gold accent under header
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.08),
        SLIDE_WIDTH, Inches(0.07),
        fill=GOLD
    )
    # Title
    add_text_box(
        slide, Inches(0.6), Inches(0.25),
        Inches(12), Inches(0.65),
        title,
        FONT_TITLE, Pt(32), WHITE, bold=True
    )
    if subtitle:
        add_text_box(
            slide, Inches(0.6), Inches(0.82),
            Inches(12), Inches(0.3),
            subtitle,
            FONT_BODY, Pt(14), GOLD_LIGHT
        )


def make_cover_slide(slide):
    # Background: soft emerald gradient-ish via two rectangles
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT,
        fill=SOFT_BG
    )
    # Decorative emerald panel on left
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.25), SLIDE_HEIGHT,
        fill=EMERALD
    )
    # Gold line accents
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(2.4),
        Inches(1.2), Inches(0.08),
        fill=GOLD
    )
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0.35), Inches(4.7),
        Inches(1.2), Inches(0.08),
        fill=GOLD
    )
    # Main title
    add_text_box(
        slide, Inches(1.8), Inches(2.1),
        Inches(10.5), Inches(1.1),
        "TartilPro x SD Khadijah 3",
        FONT_TITLE, Pt(48), EMERALD_DARK, bold=True
    )
    # Subtitle
    add_text_box(
        slide, Inches(1.8), Inches(3.25),
        Inches(10.5), Inches(0.9),
        "Transformasi Digital Pembelajaran & Evaluasi\nTartil Al-Qur'an Berbasis Cloud",
        FONT_TITLE, Pt(24), MUTED_TEXT
    )
    # Tagline pill
    pill = add_full_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.8), Inches(4.3),
        Inches(4.2), Inches(0.55),
        fill=EMERALD
    )
    add_text_box(
        slide, Inches(1.8), Inches(4.33),
        Inches(4.2), Inches(0.55),
        "Modern Islamic Tech Solution",
        FONT_BODY, Pt(16), WHITE, bold=True, align=PP_ALIGN.CENTER
    )
    # Islamic geometric-ish decorative circles
    add_full_shape(
        slide, MSO_SHAPE.OVAL,
        SLIDE_WIDTH - Inches(3.2), Inches(0.8),
        Inches(1.6), Inches(1.6),
        fill=GOLD_LIGHT
    )
    add_full_shape(
        slide, MSO_SHAPE.OVAL,
        SLIDE_WIDTH - Inches(2.4), Inches(1.9),
        Inches(1.1), Inches(1.1),
        fill=EMERALD
    )
    add_text_box(
        slide, SLIDE_WIDTH - Inches(2.4), Inches(2.1),
        Inches(1.1), Inches(0.8),
        "☪",
        FONT_TITLE, Pt(36), WHITE, bold=True, align=PP_ALIGN.CENTER
    )
    # Bottom text
    add_text_box(
        slide, Inches(1.8), Inches(5.6),
        Inches(10), Inches(0.6),
        "Slide Deck Presentasi Produk | 2026",
        FONT_BODY, Pt(16), MUTED_TEXT
    )


def make_bgsolusi_slide(slide):
    add_header_bar(slide, "Latar Belakang & Solusi", subtitle="Mengapa SD Khadijah 3 butuh transformasi digital?")

    # Left card: Latar Belakang
    card1 = add_full_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.55),
        Inches(5.9), Inches(5.0),
        fill=WHITE
    )
    card1.shadow.inherit = False
    # Card header
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.55),
        Inches(5.9), Inches(0.7),
        fill=EMERALD
    )
    add_text_box(
        slide, Inches(0.65), Inches(1.65),
        Inches(5.5), Inches(0.6),
        "Tantangan Saat Ini",
        FONT_TITLE, Pt(20), WHITE, bold=True
    )
    bullets1 = [
        "Pencatatan hafalan & bacaan masih manual di buku/jurnal.",
        "Rekap nilai tartil memakan waktu berhari-hari.",
        "Orang tua terbatas mengakses perkembangan anak.",
        "Data rapor rentan hilang dan sulit dianalisis."
    ]
    add_bullet_list(
        slide, Inches(0.75), Inches(2.45),
        Inches(5.4), Inches(4.0),
        bullets1, FONT_BODY, Pt(15), DARK_TEXT
    )

    # Right card: Solusi
    card2 = add_full_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.9), Inches(1.55),
        Inches(5.9), Inches(5.0),
        fill=WHITE
    )
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(6.9), Inches(1.55),
        Inches(5.9), Inches(0.7),
        fill=GOLD
    )
    add_text_box(
        slide, Inches(7.05), Inches(1.65),
        Inches(5.5), Inches(0.6),
        "Solusi TartilPro",
        FONT_TITLE, Pt(20), EMERALD_DARK, bold=True
    )
    bullets2 = [
        "Digitalisasi pencatatan hafalan & bacaan Al-Qur'an real-time.",
        "Penilaian tartil terstruktur: tajwid, makhraj, hafalan.",
        "Portal wali murid untuk transparansi perkembangan anak.",
        "Dashboard analytics untuk pengambilan keputusan cepat."
    ]
    add_bullet_list(
        slide, Inches(7.15), Inches(2.45),
        Inches(5.4), Inches(4.0),
        bullets2, FONT_BODY, Pt(15), DARK_TEXT
    )

    add_footer(slide)


def make_modul_slide(slide, title, subtitle, bullets, icon, accent=EMERALD):
    add_header_bar(slide, "Modul Utama Sistem", subtitle=subtitle)

    # Large icon circle
    circle = add_full_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(0.9), Inches(1.95),
        Inches(1.8), Inches(1.8),
        fill=GOLD_LIGHT
    )
    add_text_box(
        slide, Inches(0.9), Inches(2.15),
        Inches(1.8), Inches(1.5),
        icon,
        FONT_TITLE, Pt(48), accent, bold=True, align=PP_ALIGN.CENTER
    )

    # Module title
    add_text_box(
        slide, Inches(3.0), Inches(2.1),
        Inches(9.5), Inches(0.7),
        title,
        FONT_TITLE, Pt(28), EMERALD_DARK, bold=True
    )
    # Gold underline
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(3.0), Inches(2.75),
        Inches(2.0), Inches(0.07),
        fill=GOLD
    )

    # Bullet list
    add_bullet_list(
        slide, Inches(3.0), Inches(3.05),
        Inches(9.5), Inches(3.5),
        bullets, FONT_BODY, Pt(18), DARK_TEXT
    )

    add_footer(slide)


def make_fitur_unggulan_slide(slide):
    add_header_bar(slide, "Fitur Unggulan", subtitle="Inovasi yang memudahkan guru, siswa, dan wali murid")

    features = [
        ("Live Scoring via Mobile", "Guru dapat memberi nilai langsung dari ponsel saat pembelajaran berlangsung.", "📱"),
        ("Rapor Tartil Digital Otomatis", "Rapor nilai tartil, absensi, dan hafalan tergenerasi otomatis per semester.", "📄"),
        ("Rekap Presensi Pembiasaan", "Presensi harian dan pembiasaan bacaan terekap tanpa input ganda.", "✅"),
    ]

    card_w = Inches(12.0)
    card_h = Inches(1.45)
    start_y = Inches(1.55)
    gap = Inches(0.28)

    for i, (ftitle, fdesc, ficon) in enumerate(features):
        y = start_y + i * (card_h + gap)
        x = Inches(0.65)

        # Card background
        add_full_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y,
            card_w, card_h,
            fill=WHITE
        )

        # Left gold accent strip
        add_full_shape(
            slide, MSO_SHAPE.RECTANGLE,
            x, y,
            Inches(0.12), card_h,
            fill=GOLD
        )

        # Icon circle
        icon_size = Inches(0.95)
        icon_x = x + Inches(0.35)
        icon_y = y + (card_h - icon_size) / 2
        add_full_shape(
            slide, MSO_SHAPE.OVAL,
            icon_x, icon_y,
            icon_size, icon_size,
            fill=GOLD_LIGHT
        )
        add_text_box(
            slide, icon_x, icon_y + Inches(0.05),
            icon_size, Inches(0.9),
            ficon,
            FONT_TITLE, Pt(32), GOLD, bold=True, align=PP_ALIGN.CENTER
        )

        # Title
        title_x = icon_x + icon_size + Inches(0.35)
        add_text_box(
            slide, title_x, y + Inches(0.15),
            Inches(9.8), Inches(0.45),
            ftitle,
            FONT_TITLE, Pt(20), EMERALD_DARK, bold=True
        )

        # Description
        add_text_box(
            slide, title_x, y + Inches(0.55),
            Inches(9.8), Inches(0.75),
            fdesc,
            FONT_BODY, Pt(15), MUTED_TEXT
        )

    add_footer(slide)


def make_server_slide(slide):
    add_header_bar(slide, "Keunggulan Server (Laravel Cloud)", subtitle="Infrastruktur andal untuk skala sekolah")

    servers = [
        ("Autoscaling", "Bebas down saat ujian dan traffic tinggi; sumber daya menyesuaikan beban.", "⚡"),
        ("Keamanan & Backup", "Proteksi data siswa terjamin dengan backup otomatis berkala.", "🔒"),
        ("Performa Ultra Cepat", "Server cloud Laravel dioptimalkan untuk respons aplikasi real-time.", "🚀"),
        ("Zero Maintenance", "Fokus mengajar; server, patch, dan monitoring ditangani tim kami.", "🛡"),
    ]

    # Two columns
    col1 = servers[:2]
    col2 = servers[2:]

    for col_idx, col_data in enumerate([col1, col2]):
        base_x = Inches(0.7 + col_idx * 6.0)
        for row_idx, (title, desc, icon) in enumerate(col_data):
            y = Inches(1.7 + row_idx * 2.45)
            # Card
            add_full_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                base_x, y,
                Inches(5.8), Inches(2.2),
                fill=WHITE
            )
            # Icon circle
            add_full_shape(
                slide, MSO_SHAPE.OVAL,
                base_x + Inches(0.25), y + Inches(0.45),
                Inches(1.1), Inches(1.1),
                fill=GOLD_LIGHT
            )
            add_text_box(
                slide, base_x + Inches(0.25), y + Inches(0.55),
                Inches(1.1), Inches(0.95),
                icon,
                FONT_TITLE, Pt(28), GOLD, bold=True, align=PP_ALIGN.CENTER
            )
            add_text_box(
                slide, base_x + Inches(1.55), y + Inches(0.35),
                Inches(4.0), Inches(0.5),
                title,
                FONT_TITLE, Pt(18), EMERALD_DARK, bold=True
            )
            add_text_box(
                slide, base_x + Inches(1.55), y + Inches(0.95),
                Inches(4.0), Inches(1.0),
                desc,
                FONT_BODY, Pt(14), MUTED_TEXT
            )

    add_footer(slide)


def make_dampak_roadmap_slide(slide):
    add_header_bar(slide, "Dampak & Roadmap", subtitle="Manfaat nyata dan rencana go-live")

    # Left: Dampak
    add_full_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.55),
        Inches(6.0), Inches(5.0),
        fill=WHITE
    )
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.55),
        Inches(6.0), Inches(0.7),
        fill=EMERALD
    )
    add_text_box(
        slide, Inches(0.65), Inches(1.65),
        Inches(5.7), Inches(0.6),
        "Dampak yang Diharapkan",
        FONT_TITLE, Pt(20), WHITE, bold=True
    )
    impact_items = [
        "Efisiensi rekap nilai tartil naik hingga 80%.",
        "Transparansi perkembangan siswa ke orang tua meningkat.",
        "Data hafalan dan rapor tersimpan aman & terpusat.",
        "Guru lebih fokus mengajar, bukan mengelola administrasi."
    ]
    add_bullet_list(
        slide, Inches(0.75), Inches(2.45),
        Inches(5.5), Inches(4.0),
        impact_items, FONT_BODY, Pt(15), DARK_TEXT
    )

    # Right: Roadmap
    add_full_shape(
        slide, MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.55),
        Inches(6.0), Inches(5.0),
        fill=WHITE
    )
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(6.8), Inches(1.55),
        Inches(6.0), Inches(0.7),
        fill=GOLD
    )
    add_text_box(
        slide, Inches(6.95), Inches(1.65),
        Inches(5.7), Inches(0.6),
        "Roadmap Go-Live",
        FONT_TITLE, Pt(20), EMERALD_DARK, bold=True
    )
    roadmap_items = [
        "Minggu 1-2: Data awal siswa, guru, dan kelas tersedia.",
        "Minggu 3: Pelatihan penggunaan aplikasi untuk guru & admin.",
        "Minggu 4: Uji coba parallel dengan pencatatan manual.",
        "Bulan 2: Go-live penuh & evaluasi bersama SD Khadijah 3."
    ]
    add_bullet_list(
        slide, Inches(7.0), Inches(2.45),
        Inches(5.5), Inches(4.0),
        roadmap_items, FONT_BODY, Pt(15), DARK_TEXT
    )

    add_footer(slide)


def make_closing_slide(slide):
    add_full_shape(
        slide, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT,
        fill=EMERALD_DARK
    )
    # Decorative gold rings
    add_full_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(4.8), Inches(0.8),
        Inches(3.6), Inches(3.6),
        fill=GOLD
    )
    add_full_shape(
        slide, MSO_SHAPE.OVAL,
        Inches(5.15), Inches(1.15),
        Inches(2.9), Inches(2.9),
        fill=EMERALD_DARK
    )
    add_text_box(
        slide, Inches(4.8), Inches(1.7),
        Inches(3.6), Inches(1.0),
        "☪",
        FONT_TITLE, Pt(60), GOLD, bold=True, align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(0), Inches(4.6),
        SLIDE_WIDTH, Inches(1.0),
        "Terima Kasih",
        FONT_TITLE, Pt(50), WHITE, bold=True, align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(0), Inches(5.6),
        SLIDE_WIDTH, Inches(0.7),
        "TartilPro x SD Khadijah 3",
        FONT_TITLE, Pt(24), GOLD, bold=True, align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, Inches(0), Inches(6.3),
        SLIDE_WIDTH, Inches(0.5),
        "Transformasi Digital Pembelajaran & Evaluasi Tartil Al-Qur'an",
        FONT_BODY, Pt(14), WHITE, align=PP_ALIGN.CENTER
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # 1. Cover
    make_cover_slide(prs.slides.add_slide(blank_layout))

    # 2. Latar Belakang & Solusi
    make_bgsolusi_slide(prs.slides.add_slide(blank_layout))

    # 3-6. Modul Utama
    modul_data = [
        (
            "Modul Manajemen Siswa & Halaqah",
            "Kelola data siswa, kelas, dan halaqah tartil dengan terstruktur",
            [
                "Pendaftaran dan pengelompokan siswa per kelas/halaqah.",
                "Manajemen guru pembimbing dan pembagian kelas tartil.",
                "Tracking perpindahan, kenaikan kelas, dan status siswa.",
                "Integrasi data dengan semester & tahun ajaran aktif."
            ],
            "👥"
        ),
        (
            "Modul Penilaian & Tartil Tracking",
            "Penilaian komprehensif: tajwid, makhraj, dan hafalan",
            [
                "Input nilai B/C/K untuk setiap siswa secara batch harian.",
                "Penilaian aspek tajwid, makhraj, dan kelancapan bacaan.",
                "Tracking hafalan per juz 1-30 dengan persentase kumulatif.",
                "Rekap R2 (rata-rata akhir) otomatis untuk rapor."
            ],
            "📊"
        ),
        (
            "Modul Portal Wali Murid & Lapor Digital",
            "Akses transparan orang tua ke perkembangan anak",
            [
                "Login aman untuk wali murid menggunakan NIS dan nomor HP.",
                "Lihat nilai tartil, absensi, dan hafalan anak secara real-time.",
                "Unduh rapor tartil digital kapan saja.",
                "Notifikasi perkembangan dan laporan mingguan."
            ],
            "👨‍👩‍👧"
        ),
        (
            "Modul Analytics & Executive Dashboard",
            "Wawasan data untuk kepala sekolah dan admin",
            [
                "Dashboard rekap nilai, absensi, dan hafalan per kelas.",
                "Visualisasi progress siswa dan perbandingan antar semester.",
                "Export laporan PDF & Excel untuk arsip dan rapat.",
                "Audit trail lengkap atas setiap perubahan data."
            ],
            "📈"
        ),
    ]
    for title, subtitle, bullets, icon in modul_data:
        make_modul_slide(prs.slides.add_slide(blank_layout), title, subtitle, bullets, icon)

    # 7. Fitur Unggulan
    make_fitur_unggulan_slide(prs.slides.add_slide(blank_layout))

    # 8. Keunggulan Server
    make_server_slide(prs.slides.add_slide(blank_layout))

    # 9. Dampak & Roadmap
    make_dampak_roadmap_slide(prs.slides.add_slide(blank_layout))

    # 10. Closing
    make_closing_slide(prs.slides.add_slide(blank_layout))

    output_path = "TartilPro_SD_Khadijah3_SlideDeck.pptx"
    prs.save(output_path)
    print(f"Slide deck berhasil dibuat: {output_path}")


if __name__ == "__main__":
    main()
